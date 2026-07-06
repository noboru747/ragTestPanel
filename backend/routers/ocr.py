import asyncio
import base64
import os
import tempfile
from fastapi import APIRouter, UploadFile, File, HTTPException
from services.ollama_service import ollama

router = APIRouter()

SUPPORTED_IMAGE_TYPES = {"image/jpeg", "image/png", "image/webp", "image/gif"}


def _guess_type_by_ext(filename: str) -> str:
    ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
    return {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "doc":  "application/msword",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "xls":  "application/vnd.ms-excel",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ppt":  "application/vnd.ms-powerpoint",
        "jpg":  "image/jpeg",
        "jpeg": "image/jpeg",
        "png":  "image/png",
        "webp": "image/webp",
        "gif":  "image/gif",
    }.get(ext, "")


@router.post("/extract")
async def extract_text(file: UploadFile = File(...)):
    content = await file.read()
    fname = file.filename or ""
    # browsers sometimes send application/octet-stream — fall back to extension
    content_type = file.content_type or ""
    if not content_type or content_type == "application/octet-stream":
        content_type = _guess_type_by_ext(fname)

    if content_type in SUPPORTED_IMAGE_TYPES:
        b64 = base64.b64encode(content).decode()
        try:
            text = await ollama.vision_ocr(b64, filename=fname)
            return {"filename": fname, "text": text, "method": "vision_ocr"}
        except Exception:
            pass
        # fallback: tesseract OCR
        try:
            text = _extract_image_tesseract(content)
            return {"filename": fname, "text": text, "method": "tesseract"}
        except Exception as exc:
            raise HTTPException(
                status_code=503,
                detail=f"圖片 OCR 失敗（vision 模型未就緒，tesseract 也無法處理）。錯誤：{exc}",
            )

    if content_type == "application/pdf":
        text = _extract_pdf(content)
        return {"filename": fname, "text": text, "method": "pdf_parse"}

    if "wordprocessingml" in content_type:
        text = _extract_docx(content)
        return {"filename": fname, "text": text, "method": "docx_parse"}

    if content_type == "application/msword" or fname.lower().endswith(".doc"):
        text = await _extract_doc(content, fname)
        return {"filename": fname, "text": text, "method": "antiword"}

    if "spreadsheetml" in content_type:
        text = _extract_xlsx(content)
        return {"filename": fname, "text": text, "method": "xlsx_parse"}

    if "presentationml" in content_type:
        text = _extract_pptx(content)
        return {"filename": fname, "text": text, "method": "pptx_parse"}

    raise HTTPException(status_code=415, detail=f"不支援的檔案類型：{content_type or fname}")


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader
    import io
    reader = PdfReader(io.BytesIO(content))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(content: bytes) -> str:
    from docx import Document
    import io
    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook
    import io
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"=== {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            row_str = "\t".join(str(c) if c is not None else "" for c in row)
            if row_str.strip():
                lines.append(row_str)
    return "\n".join(lines)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    import io
    prs = Presentation(io.BytesIO(content))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        lines.append(text)
    return "\n".join(lines)


def _extract_image_tesseract(content: bytes) -> str:
    import pytesseract
    from PIL import Image
    import io
    img = Image.open(io.BytesIO(content))
    return pytesseract.image_to_string(img, lang="chi_tra+chi_sim+eng")


async def _extract_doc(content: bytes, filename: str) -> str:
    tmp = tempfile.NamedTemporaryFile(suffix=".doc", delete=False)
    try:
        tmp.write(content)
        tmp.close()
        proc = await asyncio.create_subprocess_exec(
            "antiword", tmp.name,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=30)
        if proc.returncode != 0:
            err = stderr.decode(errors="replace").strip()
            raise HTTPException(
                status_code=415,
                detail=f"無法解析 .doc 檔案（{err or '格式不支援'}）。請在 Word 中另存為 .docx 後重新上傳。",
            )
        return stdout.decode(errors="replace")
    finally:
        os.unlink(tmp.name)
